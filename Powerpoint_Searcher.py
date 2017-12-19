from pptx import Presentation
from os import listdir
import os
from tkinter import *
from tkinter import ttk
import tkinter as tk
from tkinter import filedialog
import tkinter.scrolledtext as tkst
import win32com.client  

directoryPath = ''
text_runs = []

print("loading...")

def findWord(filePath):
    search = searchWord.get()
    prs = Presentation(filePath)
    found = 0
    foundInFile = []
    foundSlides = []
    
    for count, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:             
                    if search.lower() in run.text.lower():
                        foundSlides.append(count)
                        found = 1

    if found == 1:
        foundInFile.append([os.path.basename(os.path.dirname(filePath)), os.path.basename(filePath), filePath])    
        foundInFile.append(foundSlides)
        text_runs.append(foundInFile)


def findFiles(path):
    #path = r'C:\Users\owner\Dropbox\Python\Projects\pptx_folder'
    all_files = []

    checkPpt(path)

    for dir, subdir, files in os.walk(path):
        
        for file in files:
            filename, file_ext = os.path.splitext(file)

            print("working")
            if "~$" in filename:                             # account for currently opened files; the unopened version is still accounted for thus still searched through   
                continue

            if file_ext == ".pptx":               
                all_files.append((os.path.join(dir, file)))
                
        
        if len(subdir) > 0 and subDirectory.get():
            del subdir[:]

    return all_files

# python-pptx library cannot inherently process .ppt ext files, need to convert to .pptx
def checkPpt(path):
    for dir, subdir, files in os.walk(path):
        
        for file in files:
            filename, file_ext = os.path.splitext(file)
        
            if file_ext == ".ppt" and not (os.path.isfile((os.path.join(dir, filename))+".pptx")):           
                convert((os.path.join(dir, file)), (os.path.join(dir, filename)))

# convert from .ppt ext to .pptx
def convert(convertFullPath, filePath):                     
    Application = win32com.client.Dispatch("PowerPoint.Application")
    Application.Visible = True
    Presentation = Application.Presentations.Open(convertFullPath)
    Presentation.Saveas(filePath + ".pptx")
    Presentation.Close()
    Application.Quit()

def main():
    global textBox
    textBox.config(state=NORMAL)    
    text_runs.clear()
    for i in findFiles(directoryPath):
        findWord(i)

    output(text_runs)


root = tk.Tk()
root.title("Powerpoint Searcher")

showDir = StringVar()
showDir.set("Current Folder: "+directoryPath)

subDirectory = tk.IntVar()
subDirectory.set(0)

def open():
    global directoryPath
    global showDir
    directoryPath =  filedialog.askdirectory()
    directoryPath = (os.path.abspath(directoryPath))
    showDir.set("Current Folder: "+directoryPath)

def output(text_runs):
    global textBox
    textBox.delete('1.0', END)
    for count, found in enumerate(text_runs):
        textBox.insert('1.0',"SLIDES: " + str(found[1]) + " "*5 + " FILE: " + found[0][1] + " "*5 + " FOLDER: " + found[0][0], count)
        textBox.insert('1.0', '\n')
        textBox.tag_config(count, foreground="blue", font="calibri")
        textBox.tag_bind(count, "<ButtonRelease-1>", lambda event, i=found: startApp(Event, i))
        textBox.update_idletasks()
     
    textBox.config(state=DISABLED)

def startApp(event, found):
    os.startfile(found[0][2])



mainframe = ttk.Frame(root, padding="23 23 42 42")
mainframe.grid(column=0, row=0)

searchWord= ttk.Entry(mainframe, width=100)
searchWord.grid(column=1, row=5, sticky=(W, E))

ttk.Button(mainframe, text="Search", command = main, width=20).grid(column=2, row=5)
ttk.Button(mainframe, text="Open Folder", command = open, width=20).grid(column=2, row=8)

ttk.Checkbutton(mainframe, text="Do not include folders in folder", onvalue=1, offvalue=0, variable=subDirectory).grid(column=2, row=9)

ttk.Label(mainframe, textvariable=showDir).grid(column=1, row=7, sticky=(W, E))

textBox = tkst.ScrolledText(mainframe, cursor="arrow", width = 100)
textBox.grid(column=1, row=10, columnspan=2)
textBox.config(state=DISABLED)    

root.mainloop()






















'''
TODO:
access within directories of directories

corner cases for search:
    ie. spelled it wrong, embeded in a chart?

UI: Tkinter has open directory, check other UI libraries 

Testing:
    500 files at once 


'''
